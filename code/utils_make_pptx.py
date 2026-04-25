from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x2B, 0x55)   # slide background / header fill
MID_BLUE  = RGBColor(0x1A, 0x5F, 0xA8)   # accent bar / bullet marker
GOLD      = RGBColor(0xFF, 0xB7, 0x00)   # highlight / title accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF2, 0xF6, 0xFC)   # content-slide background
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MED_TEXT  = RGBColor(0x2C, 0x3E, 0x6B)
GRAY_LINE = RGBColor(0xCC, 0xD6, 0xE8)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background() if line_rgb is None else None
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_paragraph(tf, text, font_size=16, bold=False, italic=False,
                  color=DARK_TEXT, align=PP_ALIGN.LEFT,
                  space_before=6, indent_level=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def slide_chrome(slide, title_text, subtitle_text=""):
    """Standard header bar + optional subtitle for content slides."""
    # light background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
    # navy header bar
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=NAVY)
    # gold accent strip under header
    add_rect(slide, 0, 1.15, 13.33, 0.07, fill_rgb=GOLD)
    # title text
    add_textbox(slide, title_text, 0.35, 0.18, 12.6, 0.85,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text, 0.35, 0.88, 12.6, 0.35,
                    font_size=14, italic=True, color=GOLD, align=PP_ALIGN.LEFT)
    # slide number placeholder (bottom right)
    return slide


def content_box(slide, left, top, width, height,
                header=None, header_color=MID_BLUE):
    """A rounded-ish white card with optional colored header bar."""
    add_rect(slide, left, top, width, height, fill_rgb=WHITE,
             line_rgb=GRAY_LINE, line_width_pt=1)
    if header:
        add_rect(slide, left, top, width, 0.38, fill_rgb=header_color)
        add_textbox(slide, header, left + 0.12, top + 0.05, width - 0.2, 0.32,
                    font_size=15, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# full navy background
add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
# diagonal accent block (top-right)
add_rect(s1, 9.5, 0, 3.83, 7.5, fill_rgb=MID_BLUE)
# gold bar
add_rect(s1, 0, 5.75, 13.33, 0.09, fill_rgb=GOLD)

# USC Trojan decorative band
add_rect(s1, 0, 0, 0.18, 7.5, fill_rgb=GOLD)

# "NLP Project" tag
add_textbox(s1, "NLP PROJECT  •  USC", 0.38, 0.28, 9.0, 0.45,
            font_size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# Main title
add_textbox(s1,
    "Beyond Binary Accuracy",
    0.38, 0.85, 9.0, 1.1,
    font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(s1,
    "Evaluating Long-form Context & Challenging Cases\nin Movie Review Sentiment Analysis",
    0.38, 1.95, 8.8, 1.1,
    font_size=22, bold=False, color=RGBColor(0xB0, 0xC8, 0xF0), align=PP_ALIGN.LEFT)

# Separator
add_rect(s1, 0.38, 3.15, 5.2, 0.055, fill_rgb=GOLD)

# Author block
authors = [
    ("Raghav Sarmukaddam", "Jimmy Taravia", "Dhairya Shah"),
    ("Yash Desai", "Vivek Lakhani", ""),
]
y = 3.32
for row in authors:
    line = "   •   ".join([a for a in row if a])
    add_textbox(s1, line, 0.38, y, 8.8, 0.42,
                font_size=15, color=WHITE, align=PP_ALIGN.LEFT)
    y += 0.42

add_textbox(s1, "University of Southern California", 0.38, y + 0.1, 8.8, 0.4,
            font_size=14, italic=True, color=RGBColor(0xB0, 0xC8, 0xF0))

# Right panel decoration
add_textbox(s1, "🎬", 10.1, 1.8, 2.5, 2.5, font_size=80, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "Sentiment\nAnalysis", 9.9, 4.2, 3.0, 1.2,
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation & Research Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_chrome(s2, "Motivation & Research Problem")

# Left column — problem statement card
content_box(s2, 0.35, 1.35, 6.1, 5.7, "The Problem", MID_BLUE)

txb = s2.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(5.8), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = "High benchmark accuracy ≠ robust understanding"
r0.font.size = Pt(17)
r0.font.bold = True
r0.font.color.rgb = NAVY
r0.font.name = "Calibri"

bullets = [
    "Standard datasets are dominated by easy examples — models learn surface keywords",
    "Difficult reviews rely on long-form context, negation, and mixed sentiment",
    "A model scoring 89% overall may still fail systematically on hard cases",
]
for b in bullets:
    add_paragraph(tf, "▸  " + b, font_size=15, color=DARK_TEXT, space_before=10)

add_paragraph(tf, "", font_size=8, space_before=4)
add_paragraph(tf, "Core Question:", font_size=16, bold=True, color=MID_BLUE, space_before=6)
add_paragraph(tf,
    "Does the progression from classical ML → deep learning → LLMs yield "
    "genuine improvements on difficult movie reviews?",
    font_size=15, italic=True, color=NAVY, space_before=6)

# Right column — example + insight
content_box(s2, 6.75, 1.35, 6.2, 2.55, "Why Movie Reviews?", GOLD)

txb2 = s2.shapes.add_textbox(Inches(6.9), Inches(1.85), Inches(5.9), Inches(2.0))
txb2.word_wrap = True
tf2 = txb2.text_frame; tf2.word_wrap = True
items = [
    "Long multi-paragraph narratives",
    "Sentiment progression & reversal",
    "Rhetorical contrast & nuanced judgment",
    "Implicit polarity (sarcasm, irony)",
]
p_first = tf2.paragraphs[0]
p_first.alignment = PP_ALIGN.LEFT
r_first = p_first.add_run(); r_first.text = items[0]
r_first.font.size = Pt(15); r_first.font.color.rgb = DARK_TEXT; r_first.font.name = "Calibri"
for item in items[1:]:
    add_paragraph(tf2, item, font_size=15, color=DARK_TEXT, space_before=7)

# Example quote box
content_box(s2, 6.75, 4.1, 6.2, 2.95, "Challenging Example", RGBColor(0xC0, 0x39, 0x2B))
add_textbox(s2,
    '"The plot was not bad at all, but I\ncannot say the acting was anything\nother than painfully disappointing."',
    6.9, 4.6, 5.9, 1.4,
    font_size=14, italic=True, color=DARK_TEXT)
add_textbox(s2,
    "→  Negation + mixed sentiment  →  Misleads bag-of-words models",
    6.9, 5.9, 5.9, 0.55,
    font_size=13, bold=True, color=RGBColor(0xC0, 0x39, 0x2B))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Literature Review
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_chrome(s3, "Background & Literature Review", "Evolution of Sentiment Analysis")

# Timeline arrow base
add_rect(s3, 0.35, 3.5, 12.6, 0.12, fill_rgb=GRAY_LINE)

# Four era boxes
eras = [
    ("Traditional ML", "Naïve Bayes & SVM\nwith TF-IDF features\n\nBaid 2017; Kumar 2018",
     MID_BLUE, 0.35),
    ("Deep Learning", "CNN classifiers\ncapture local n-gram\npatterns\n\nHaque et al.",
     RGBColor(0x27, 0x8E, 0x68), 3.6),
    ("Transformers", "BERT: contextualized\nrepresentations via\nself-attention\n\nPandey 2024",
     RGBColor(0x8E, 0x44, 0xAD), 6.85),
    ("LLMs", "Zero/few-shot\nclassification;\nimplicit reasoning\n\nGosai 2025",
     RGBColor(0xC0, 0x39, 0x2B), 10.1),
]

arrow_labels = ["→", "→", "→"]
for i, (title, body, color, x) in enumerate(eras):
    add_rect(s3, x, 1.4, 2.9, 2.2, fill_rgb=color)
    add_textbox(s3, title, x + 0.1, 1.45, 2.7, 0.5,
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s3, x, 1.88, 2.9, 1.52, fill_rgb=WHITE,
             line_rgb=color, line_width_pt=1.5)
    add_textbox(s3, body, x + 0.12, 1.93, 2.66, 1.45,
                font_size=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    # dot on timeline
    dot = s3.shapes.add_shape(9, Inches(x + 1.3), Inches(3.44), Inches(0.24), Inches(0.24))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    if i < 3:
        add_textbox(s3, "→", x + 3.05, 2.3, 0.45, 0.6,
                    font_size=26, bold=True, color=GRAY_LINE, align=PP_ALIGN.CENTER)

# Gap insight below timeline
content_box(s3, 0.35, 3.8, 12.6, 2.95, "Key Insight from Literature", NAVY)
add_textbox(s3,
    "Despite steady accuracy gains across these eras, open challenges remain:\n"
    "sarcasm  •  negation  •  contextual ambiguity  •  real-world robustness\n\n"
    "→  Overall test accuracy alone is insufficient to evaluate genuine sentiment understanding  (Gosai et al., 2025)",
    0.5, 4.28, 12.2, 2.2,
    font_size=16, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our Approach
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_chrome(s4, "Our Approach", "Comparative evaluation pipeline across model generations")

# Left: model suite
content_box(s4, 0.35, 1.35, 5.8, 5.75, "Model Suite", MID_BLUE)

models = [
    ("Traditional ML",   "Naïve Bayes + SVM (TF-IDF)",            "✓ Done"),
    ("Hybrid ML",        "SVM + Lexicon polarity features",        "✓ Done"),
    ("Deep Learning",    "CNN (Kim 2014, GloVe embeddings)",       "⚙ In Progress"),
    ("Transformer",      "Fine-tuned BERT",                        "⏳ Upcoming"),
    ("LLM (stretch)",    "Llama 3 / Mistral — zero-shot",          "⏳ Stretch Goal"),
]
status_colors = {
    "✓ Done":          RGBColor(0x27, 0x8E, 0x68),
    "⚙ In Progress":  RGBColor(0xFF, 0xB7, 0x00),
    "⏳ Upcoming":     MID_BLUE,
    "⏳ Stretch Goal": RGBColor(0x95, 0xA5, 0xA6),
}

y = 1.85
for tier, desc, status in models:
    add_rect(s4, 0.4, y, 5.7, 0.88,
             fill_rgb=WHITE, line_rgb=GRAY_LINE, line_width_pt=0.8)
    add_rect(s4, 0.4, y, 0.14, 0.88, fill_rgb=status_colors[status])
    add_textbox(s4, tier, 0.62, y + 0.05, 3.0, 0.38,
                font_size=14, bold=True, color=NAVY)
    add_textbox(s4, desc, 0.62, y + 0.42, 3.5, 0.38,
                font_size=12, color=MED_TEXT)
    add_textbox(s4, status, 4.05, y + 0.2, 1.9, 0.5,
                font_size=12, bold=True, color=status_colors[status], align=PP_ALIGN.RIGHT)
    y += 1.0

# Middle: datasets
content_box(s4, 6.45, 1.35, 3.3, 2.7, "Datasets", RGBColor(0x27, 0x8E, 0x68))
add_textbox(s4,
    "Primary\nIMDb Large Movie Review\n50K reviews  •  balanced\nLong-form, multi-paragraph",
    6.6, 1.83, 3.0, 1.15, font_size=13, color=DARK_TEXT)
add_rect(s4, 6.6, 3.0, 3.0, 0.05, fill_rgb=GRAY_LINE)
add_textbox(s4,
    "Secondary\nStanford Sentiment Treebank\nShorter phrase-level annotations\nDocument vs phrase-level contrast",
    6.6, 3.07, 3.0, 1.0, font_size=13, color=DARK_TEXT)

# Right: evaluation slices
content_box(s4, 6.45, 4.2, 3.3, 2.9, "Challenging Slices", RGBColor(0xC0, 0x39, 0x2B))
slices = [
    ("📏", "Long Reviews",     "> 500 tokens"),
    ("🔁", "Negation-heavy",   "not, never, hardly…"),
    ("⚖️", "Mixed Polarity",   "praise + criticism"),
    ("❌", "Hard Errors",      "post-hoc analysis"),
]
sy = 4.68
for icon, name, detail in slices:
    add_textbox(s4, icon, 6.5, sy, 0.5, 0.45, font_size=16, align=PP_ALIGN.CENTER)
    add_textbox(s4, name, 6.95, sy, 1.6, 0.28, font_size=13, bold=True, color=DARK_TEXT)
    add_textbox(s4, detail, 6.95, sy + 0.27, 2.6, 0.25, font_size=11, italic=True, color=MED_TEXT)
    sy += 0.6

# Right-most: goal box
content_box(s4, 10.05, 1.35, 2.9, 5.75, "Core Goal", NAVY)
add_textbox(s4,
    "Move beyond\noverall binary\naccuracy\n\n↓\n\nMeasure genuine\nrobustness on\nchallenging\nreview cases",
    10.18, 1.85, 2.65, 4.8,
    font_size=15, bold=False, color=DARK_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Progress: Data & Preprocessing
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_chrome(s5, "Progress: Data & Preprocessing")

# IMDb card
content_box(s5, 0.35, 1.35, 6.1, 3.3, "IMDb Large Movie Review Dataset", MID_BLUE)
items_imdb = [
    "50,000 reviews — 25K train / 25K test, perfectly balanced",
    "Primary benchmark: long multi-paragraph reviews",
    "Ideal for studying long-form context effects",
    "Preprocessing pipeline:",
    "  Strip HTML  •  Lowercase  •  Normalize punctuation",
    "  Tokenize  •  Remove stop words",
]
y = 1.9
for item in items_imdb:
    bold_it = ":" in item and not item.startswith(" ")
    add_textbox(s5, item, 0.5, y, 5.8, 0.38,
                font_size=14 if not item.startswith(" ") else 13,
                bold=bold_it,
                color=DARK_TEXT if not item.startswith(" ") else MED_TEXT)
    y += 0.41

# SST-2 card
content_box(s5, 0.35, 4.85, 6.1, 1.95, "SST-2 Secondary Benchmark", RGBColor(0x27, 0x8E, 0x68))
items_sst = [
    "Shorter sentence-level annotations (Socher et al., 2013)",
    "Phrase-level vs document-level model comparison",
    "Reveals whether models generalize across review granularity",
]
y = 5.35
for item in items_sst:
    add_textbox(s5, item, 0.5, y, 5.8, 0.35, font_size=13, color=DARK_TEXT)
    y += 0.42

# Pipeline diagram (right side)
content_box(s5, 6.75, 1.35, 6.2, 5.45, "Preprocessing Pipeline", NAVY)

steps = [
    ("Raw Review Text",  NAVY),
    ("Strip HTML Tags",  MID_BLUE),
    ("Lowercase + Normalize Punctuation",  MID_BLUE),
    ("Tokenization",     MID_BLUE),
    ("Stop Word Removal",MID_BLUE),
    ("Feature Vectors (TF-IDF / Embeddings)", RGBColor(0x27, 0x8E, 0x68)),
]
sy = 1.82
for i, (step, col) in enumerate(steps):
    add_rect(s5, 7.1, sy, 5.5, 0.52, fill_rgb=col,
             line_rgb=GRAY_LINE, line_width_pt=0.6)
    add_textbox(s5, step, 7.25, sy + 0.1, 5.2, 0.38,
                font_size=14, bold=(i == 0 or i == len(steps)-1),
                color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_textbox(s5, "↓", 9.6, sy + 0.53, 0.5, 0.32,
                    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    sy += 0.83


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Progress: Baseline Results
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_chrome(s6, "Progress: Baseline Model Results", "IMDb test set — all models use unigram+bigram TF-IDF")

# Results table (left)
content_box(s6, 0.35, 1.35, 7.0, 4.5, "Baseline Performance on IMDb", MID_BLUE)

# Table headers
col_x = [0.5, 3.55, 5.25]
col_w = [3.0, 1.6, 1.6]
header_row = ["Model", "Accuracy", "F1 Score"]

add_rect(s6, 0.4, 1.83, 6.9, 0.52, fill_rgb=NAVY)
for j, (hdr, cx, cw) in enumerate(zip(header_row, col_x, col_w)):
    add_textbox(s6, hdr, cx, 1.87, cw, 0.42,
                font_size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rows = [
    ("Naïve Bayes (TF-IDF)",   "86.2%", "86.0%", LIGHT_BG),
    ("SVM (TF-IDF)",           "89.1%", "89.0%", WHITE),
    ("SVM + Hybrid Features",  "89.5%", "89.3%", RGBColor(0xE8, 0xF5, 0xE9)),
]
ry = 2.37
for name, acc, f1, bg in rows:
    add_rect(s6, 0.4, ry, 6.9, 0.62, fill_rgb=bg, line_rgb=GRAY_LINE, line_width_pt=0.5)
    add_textbox(s6, name, col_x[0], ry + 0.12, col_w[0], 0.42, font_size=14, color=DARK_TEXT)
    add_textbox(s6, acc,  col_x[1], ry + 0.12, col_w[1], 0.42, font_size=15, bold=True,
                color=MID_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s6, f1,   col_x[2], ry + 0.12, col_w[2], 0.42, font_size=15, bold=True,
                color=MID_BLUE, align=PP_ALIGN.CENTER)
    ry += 0.64

# Key takeaways from table
content_box(s6, 0.35, 6.0, 7.0, 1.1, "Key Takeaways", NAVY)
add_textbox(s6,
    "Hybrid features (+lexicon polarity) give the best overall result  •  "
    "But all models degrade on hard slices (see next slide)",
    0.5, 6.45, 6.7, 0.55, font_size=13, color=DARK_TEXT)

# Bar chart (right) — hand-drawn with rectangles
content_box(s6, 7.65, 1.35, 5.3, 5.75, "Accuracy Comparison", RGBColor(0x27, 0x8E, 0x68))

chart_x0, chart_y_base = 7.85, 6.55
bar_max_h = 3.8   # inches for 100%
bar_w     = 0.9
gap       = 0.52

bars = [
    ("Naïve\nBayes", 0.862, MID_BLUE),
    ("SVM\nTF-IDF",  0.891, RGBColor(0x27, 0x8E, 0x68)),
    ("SVM\nHybrid",  0.895, GOLD),
]

# Y gridlines
for pct in [0.80, 0.85, 0.90, 0.95]:
    gh = bar_max_h * pct
    gy = chart_y_base - gh
    add_rect(s6, chart_x0 + 0.05, gy, 4.2, 0.018, fill_rgb=GRAY_LINE)
    add_textbox(s6, f"{int(pct*100)}%", chart_x0 - 0.35, gy - 0.12, 0.55, 0.28,
                font_size=10, color=MED_TEXT, align=PP_ALIGN.RIGHT)

for i, (label, val, col) in enumerate(bars):
    bx = chart_x0 + 0.3 + i * (bar_w + gap)
    bh = bar_max_h * val
    by = chart_y_base - bh
    add_rect(s6, bx, by, bar_w, bh, fill_rgb=col)
    add_textbox(s6, f"{val*100:.1f}%", bx, by - 0.38, bar_w, 0.32,
                font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s6, label, bx, chart_y_base + 0.05, bar_w, 0.55,
                font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# x-axis line
add_rect(s6, chart_x0 + 0.05, chart_y_base, 4.2, 0.04, fill_rgb=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Challenging Evaluation Slices
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_chrome(s7, "Challenging Evaluation Slices",
             "Where standard accuracy metrics hide real weaknesses")

# Three slice cards
slices_info = [
    ("Long Reviews", "> 500 tokens\n(top ~20% by length)",
     "2–3 point\naccuracy drop",
     "Bag-of-words models lose track of\nsentiment when it's distributed\nacross multiple paragraphs.",
     MID_BLUE),
    ("Negation-Heavy", "Cues: not, never, no,\nhardly, without…",
     "Frequent\nmisclassification",
     "\"not bad at all\" → negative\nwords in a positive context.\nSurface statistics mislead models.",
     RGBColor(0xC0, 0x39, 0x2B)),
    ("Mixed Polarity", "High co-occurrence of\npositive + negative terms",
     "Nuanced opinions\nmisclassified",
     "Reviews blending praise and\ncriticism challenge single-label\nprediction.",
     RGBColor(0x8E, 0x44, 0xAD)),
]

sx = 0.35
for title, definition, finding, explain, col in slices_info:
    # outer card
    add_rect(s7, sx, 1.35, 4.15, 5.75, fill_rgb=WHITE,
             line_rgb=col, line_width_pt=2)
    # color header
    add_rect(s7, sx, 1.35, 4.15, 0.72, fill_rgb=col)
    add_textbox(s7, title, sx + 0.12, 1.38, 3.9, 0.62,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # definition
    add_textbox(s7, "Definition", sx + 0.15, 2.17, 3.85, 0.32,
                font_size=12, bold=True, color=col)
    add_textbox(s7, definition, sx + 0.15, 2.46, 3.85, 0.7,
                font_size=13, color=DARK_TEXT)
    # finding box
    add_rect(s7, sx + 0.1, 3.3, 3.95, 0.88, fill_rgb=col)
    add_textbox(s7, "Finding: " + finding, sx + 0.2, 3.35, 3.75, 0.78,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # explanation
    add_textbox(s7, explain, sx + 0.15, 4.35, 3.85, 1.5,
                font_size=13, color=DARK_TEXT)
    sx += 4.35

# Bottom insight banner
add_rect(s7, 0.35, 7.0, 12.6, 0.38, fill_rgb=NAVY)
add_textbox(s7,
    "→  All three slices are implemented and will be applied uniformly to every model in our pipeline",
    0.5, 7.03, 12.3, 0.32,
    font_size=13, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Current Status & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
slide_chrome(s8, "Current Status & Next Steps")

# Completed column
content_box(s8, 0.35, 1.35, 5.8, 5.75, "Completed", RGBColor(0x27, 0x8E, 0x68))

done_items = [
    ("Data Pipeline",
     "IMDb + SST-2 loaded, cleaned, and split\nPreprocessing modules verified"),
    ("Baseline Models",
     "Naïve Bayes: 86.2%\nSVM (TF-IDF): 89.1%\nSVM + Hybrid: 89.5%"),
    ("Challenging Slices",
     "Long / Negation / Mixed-polarity\nsubsets implemented & validated"),
    ("Experiment Tracking",
     "Shared evaluation code, metrics\n(Acc, Precision, Recall, Macro-F1)"),
]
dy = 1.88
for title, detail in done_items:
    add_rect(s8, 0.42, dy, 5.6, 0.95, fill_rgb=WHITE,
             line_rgb=RGBColor(0x27, 0x8E, 0x68), line_width_pt=0.8)
    add_rect(s8, 0.42, dy, 0.18, 0.95, fill_rgb=RGBColor(0x27, 0x8E, 0x68))
    add_textbox(s8, "✓  " + title, 0.68, dy + 0.04, 5.05, 0.38,
                font_size=14, bold=True, color=RGBColor(0x1A, 0x6B, 0x3C))
    add_textbox(s8, detail, 0.68, dy + 0.42, 5.05, 0.55,
                font_size=12, color=MED_TEXT)
    dy += 1.12

# In Progress column
content_box(s8, 6.45, 1.35, 3.05, 2.5, "In Progress", GOLD)
add_rect(s8, 6.52, 1.88, 2.88, 1.75, fill_rgb=WHITE,
         line_rgb=GOLD, line_width_pt=0.8)
add_rect(s8, 6.52, 1.88, 0.16, 1.75, fill_rgb=GOLD)
add_textbox(s8, "⚙  CNN Classifier", 6.76, 1.93, 2.6, 0.38,
            font_size=14, bold=True, color=RGBColor(0xB7, 0x7F, 0x00))
add_textbox(s8,
    "Kim (2014) architecture\nFilters: width 3, 4, 5\nGloVe word embeddings\nMax-over-time pooling\n\nResults expected soon",
    6.76, 2.3, 2.6, 1.25, font_size=12, color=DARK_TEXT)

# Upcoming column
content_box(s8, 6.45, 4.0, 3.05, 3.1, "Upcoming", MID_BLUE)
upcoming = [
    ("BERT Fine-tuning", "Weeks 12–13\nFull IMDb comparison"),
    ("LLM Zero-shot", "Stretch goal\nLlama 3 / Mistral"),
    ("Error Analysis", "Week 14\nFinal report"),
]
uy = 4.5
for title, detail in upcoming:
    add_textbox(s8, "⏳  " + title, 6.6, uy, 2.8, 0.35,
                font_size=13, bold=True, color=MID_BLUE)
    add_textbox(s8, detail, 6.78, uy + 0.33, 2.65, 0.38,
                font_size=11, italic=True, color=MED_TEXT)
    uy += 0.88

# Risk mitigation (right)
content_box(s8, 9.75, 1.35, 3.25, 5.75, "Risk Mitigation", NAVY)
risks = [
    ("GPU Compute", "Google Colab Pro (A100)\nGradient accumulation for BERT"),
    ("Comparison Fairness", "Shared test split & metrics\nDocument all truncation"),
    ("Slice Noise", "Manual review of 50-100\nsamples per slice"),
    ("LLM Cost", "Restrict to hard slices\n~1-2K reviews if needed"),
]
ry2 = 1.88
for risk, mitigation in risks:
    add_textbox(s8, risk, 9.9, ry2, 2.95, 0.3,
                font_size=13, bold=True, color=NAVY)
    add_textbox(s8, mitigation, 9.9, ry2 + 0.29, 2.95, 0.55,
                font_size=11, color=DARK_TEXT)
    add_rect(s8, 9.9, ry2 + 0.82, 2.9, 0.03, fill_rgb=GRAY_LINE)
    ry2 += 1.0


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Summary & Q&A
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
add_rect(s9, 0, 0, 0.2, 7.5, fill_rgb=GOLD)
add_rect(s9, 0, 6.8, 13.33, 0.7, fill_rgb=MID_BLUE)

add_textbox(s9, "Summary", 0.45, 0.3, 12.5, 0.7,
            font_size=36, bold=True, color=WHITE)
add_rect(s9, 0.45, 1.02, 5.5, 0.07, fill_rgb=GOLD)

# 4 summary points
pts = [
    ("Research Problem",
     "Do newer model generations genuinely improve on hard movie reviews?"),
    ("Our Approach",
     "Comparative pipeline: NB → SVM → CNN → BERT → LLM on IMDb & SST-2"),
    ("Completed Work",
     "Preprocessing, 3 baselines (up to 89.5%), 3 challenging slice filters"),
    ("Early Finding",
     "Baselines degrade 2-3 pts on hard slices — confirming standard metrics hide weaknesses"),
]
col_colors = [MID_BLUE, RGBColor(0x27, 0x8E, 0x68), RGBColor(0x8E, 0x44, 0xAD), RGBColor(0xC0, 0x39, 0x2B)]
bx = 0.45
for (ptitle, pdesc), col in zip(pts, col_colors):
    add_rect(s9, bx, 1.2, 3.0, 4.7, fill_rgb=col)
    add_textbox(s9, ptitle, bx + 0.12, 1.28, 2.76, 0.6,
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s9, bx + 0.12, 1.85, 2.76, 0.04, fill_rgb=WHITE)
    add_textbox(s9, pdesc, bx + 0.12, 1.97, 2.76, 3.8,
                font_size=14, color=WHITE, align=PP_ALIGN.LEFT)
    bx += 3.22

# Q&A
add_textbox(s9, "Thank you  —  Questions?", 0.45, 6.05, 12.4, 0.65,
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Team names at bottom
add_textbox(s9,
    "Raghav Sarmukaddam  •  Jimmy Taravia  •  Dhairya Shah  •  Yash Desai  •  Vivek Lakhani  •  USC",
    0.45, 6.83, 12.4, 0.38,
    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/dhairyashah5/Downloads/nlp/Beyond_Binary_Accuracy_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
